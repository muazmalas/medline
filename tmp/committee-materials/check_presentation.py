import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

p = Path(r"C:\Users\Muaz\Downloads\Pharmacy\output\presentations\MedLine_Committee_Presentation.pptx")
ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

with zipfile.ZipFile(p) as z:
    slide_names = sorted(
        [n for n in z.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)],
        key=lambda n: int(re.search(r"\d+", Path(n).stem).group()),
    )
    note_names = sorted(
        [n for n in z.namelist() if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", n)],
        key=lambda n: int(re.search(r"\d+", Path(n).stem).group()),
    )
    slides = [
        " ".join(t.text or "" for t in ET.fromstring(z.read(n)).findall(".//a:t", ns))
        for n in slide_names
    ]
    notes = [
        "\n".join(t.text or "" for t in ET.fromstring(z.read(n)).findall(".//a:t", ns))
        for n in note_names
    ]

print("slides", len(slides), "notes", len(notes))
print("notes_with_sources", sum("[Sources]" in n for n in notes))
print("speaker_counts", {f"Speaker {i}": sum(f"Speaker {i}" in n for n in notes) for i in range(1, 5)})
print("visible_start_delivery", sum("Start delivery" in s for s in slides))
print("visible_picked_up", sum("picked_up" in s for s in slides))
print(
    "notes_no_start_message",
    sum(
        "no separate driver Start delivery" in n
        or "no separate driver Start delivery action" in n
        for n in notes
    ),
)
for i, slide in enumerate(slides, 1):
    print(i, slide[:120])

prs = Presentation(p)
overflows = []
for slide_index, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if (
            shape.left < 0
            or shape.top < 0
            or shape.left + shape.width > prs.slide_width
            or shape.top + shape.height > prs.slide_height
        ):
            overflows.append((slide_index, shape.name))
print("out_of_bounds_shapes", overflows)
